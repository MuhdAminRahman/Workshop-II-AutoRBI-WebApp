import openpyxl
from openpyxl.styles import PatternFill
from pptx import Presentation
from pptx.util import Inches, Pt
from app.utils.cloudinary_util import upload_to_cloudinary
from sqlalchemy.orm import Session
from app.models.file import File
from datetime import datetime
import io

class FileService:
    """
    Manages Excel and PowerPoint file generation and versioning.
    """
    
    def __init__(self, db: Session):
        self.db = db
    
    async def generate_excel_file(self, work_id: int, masterfile_path: str) -> str:
        """
        Generate updated Excel file with extracted data.
        
        Process:
        1. Load original masterfile
        2. Find empty cells matching component fields
        3. Fill with extracted data
        4. Color-code filled cells
        5. Upload to Cloudinary
        6. Save file record in DB
        """
        from app.models.work import Work
        from app.models.equipment import Equipment
        from app.models.component import Component
        
        # Load masterfile
        wb = openpyxl.load_workbook(masterfile_path)
        ws = wb.active
        
        # Get all equipment and components for this work
        equipment_list = self.db.query(Equipment).filter(
            Equipment.work_id == work_id
        ).all()
        
        # Define which columns to fill (from your excel_manager.py)
        fields_to_fill = [
            "component_name", "phase", "fluid", "material_spec", 
            "material_grade", "insulation", "design_temp", "design_pressure",
            "operating_temp", "operating_pressure"
        ]
        
        # Green fill for auto-filled cells
        green_fill = PatternFill(start_color="90EE90", end_color="90EE90", fill_type="solid")
        
        row = 2  # Assuming row 1 is headers
        for equipment in equipment_list:
            for component in equipment.components:
                # Write to columns matching fields_to_fill
                for col_idx, field in enumerate(fields_to_fill, start=2):
                    value = getattr(component, field, None)
                    if value:
                        cell = ws.cell(row=row, column=col_idx, value=value)
                        cell.fill = green_fill
                row += 1
        
        # Save to bytes
        excel_bytes = io.BytesIO()
        wb.save(excel_bytes)
        excel_bytes.seek(0)
        
        # Upload to Cloudinary
        file_url = await upload_to_cloudinary(
            file_object=excel_bytes,
            file_name=f"work_{work_id}_excel_v{self._get_next_version(work_id)}.xlsx",
            resource_type="auto"
        )
        
        # Record in DB
        version = self._get_next_version(work_id)
        file_record = File(
            work_id=work_id,
            file_type="excel",
            version_number=version,
            file_url=file_url,
            created_by=self.current_user_id
        )
        self.db.add(file_record)
        self.db.commit()
        
        return file_url
    
    async def generate_powerpoint(self, work_id: int, ppt_template_path: str) -> str:
        """
        Generate PowerPoint from template and extracted data.
        
        Process:
        1. Load PPT template
        2. For each slide, fill in equipment/component data
        3. Add charts/tables if needed
        4. Upload to Cloudinary
        5. Save file record in DB
        """
        from app.models.equipment import Equipment
        
        # Load template
        prs = Presentation(ppt_template_path)
        
        # Get equipment data
        equipment_list = self.db.query(Equipment).filter(
            Equipment.work_id == work_id
        ).all()
        
        # Customize each slide (simplified version)
        # In real implementation, this depends on your template structure
        for idx, equipment in enumerate(equipment_list):
            if idx < len(prs.slides):
                slide = prs.slides[idx]
                
                # Find title placeholder
                for shape in slide.shapes:
                    if hasattr(shape, "text") and "Equipment" in shape.text:
                        shape.text = f"Equipment: {equipment.equipment_number}"
                    
                    # Add component table if shape is a table
                    if shape.has_table:
                        table = shape.table
                        for comp_idx, component in enumerate(equipment.components):
                            if comp_idx < len(table.rows):
                                row = table.rows[comp_idx]
                                row.cells[0].text = component.component_name
                                row.cells[1].text = component.fluid or ""
                                row.cells[2].text = component.material_spec or ""
        
        # Save to bytes
        ppt_bytes = io.BytesIO()
        prs.save(ppt_bytes)
        ppt_bytes.seek(0)
        
        # Upload to Cloudinary
        file_url = await upload_to_cloudinary(
            file_object=ppt_bytes,
            file_name=f"work_{work_id}_ppt_v{self._get_next_version(work_id)}.pptx",
            resource_type="auto"
        )
        
        # Record in DB
        version = self._get_next_version(work_id)
        file_record = File(
            work_id=work_id,
            file_type="powerpoint",
            version_number=version,
            file_url=file_url,
            created_by=self.current_user_id
        )
        self.db.add(file_record)
        self.db.commit()
        
        return file_url
    
    def _get_next_version(self, work_id: int) -> int:
        """Get next version number for files."""
        latest = self.db.query(File).filter(
            File.work_id == work_id
        ).order_by(File.version_number.desc()).first()
        
        return (latest.version_number + 1) if latest else 1