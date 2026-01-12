"""
Reports Service - FINAL VERSION
Generate Excel and PowerPoint reports from extracted equipment data
Templates downloaded from Cloudinary, reports uploaded back to Cloudinary
"""

import logging
import os
import tempfile
import httpx
from typing import Dict, List, Optional
from datetime import datetime

from sqlalchemy.orm import Session
from openpyxl import load_workbook
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

from app.models.equipment import Equipment
from app.models.component import Component
from app.utils.cloudinary_util import upload_excel_to_cloudinary, upload_ppt_to_cloudinary

logger = logging.getLogger(__name__)


# ============================================================================
# DATA MODELS FOR EXCEL/PPT
# ============================================================================

class ComponentData:
    """In-memory component data structure"""
    def __init__(self, component: Component):
        self.component_name = component.component_name
        self.phase = component.phase
        self.row_index = None
        self.data = {
            'fluid': component.fluid,
            'material_type': 'Carbon Steel' if 'SA-516' in (component.material_spec or '') else 'Stainless Steel',
            'spec': component.material_spec,
            'grade': component.material_grade,
            'insulation': component.insulation,
            'design_temp': component.design_temp,
            'design_pressure': component.design_pressure,
            'operating_temp': component.operating_temp,
            'operating_pressure': component.operating_pressure,
        }


class EquipmentData:
    """In-memory equipment data structure"""
    def __init__(self, equipment: Equipment):
        self.equipment_number = equipment.equipment_number
        self.pmt_number = equipment.pmt_number
        self.description = equipment.description
        self.components: List[ComponentData] = []


# ============================================================================
# EXCEL GENERATION
# ============================================================================

class ExcelReportGenerator:
    """Generate Excel reports from extracted data"""
    
    def __init__(self, template_path: str):
        """Initialize with Excel template path"""
        self.template_path = template_path
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"Excel template not found: {template_path}")
    
    def generate_from_equipment(self, equipment_list: List[Equipment]) -> bytes:
        """
        Generate Excel report from equipment list.
        
        Args:
            equipment_list: List of Equipment objects from database
        
        Returns:
            Excel file bytes
        """
        try:
            logger.info(f"Generating Excel report for {len(equipment_list)} equipment")
            
            # Load template
            wb = load_workbook(self.template_path)
            ws = wb['Masterfile']
            
            # Build equipment data map
            equipment_map = {}
            for equipment in equipment_list:
                equip_data = EquipmentData(equipment)
                for component in equipment.components:
                    comp_data = ComponentData(component)
                    equip_data.components.append(comp_data)
                equipment_map[equipment.equipment_number] = equip_data
            
            # Find row indices for each component in template
            self._map_component_rows(ws, equipment_map)
            
            # Fill Excel data
            self._fill_excel_data(ws, equipment_map)
            
            # Save to bytes
            import io
            output = io.BytesIO()
            wb.save(output)
            output.seek(0)
            
            logger.info("✅ Excel report generated successfully")
            return output.getvalue()
        
        except Exception as e:
            logger.error(f"❌ Error generating Excel: {str(e)}")
            raise
    
    def _map_component_rows(self, ws, equipment_map: Dict[str, EquipmentData]):
        """
        Map component row indices from template.
        Reads the Masterfile sheet to find which row each component is in.
        """
        try:
            current_row = 7  # Start from row 7 (after headers)
            current_equipment = None
            
            while current_row <= ws.max_row and current_row <= 100:
                equipment_number = self._get_cell_value(ws, f'B{current_row}')
                component_name = self._get_cell_value(ws, f'E{current_row}')
                
                # New equipment found
                if equipment_number and equipment_number not in ['EQUIPMENT NO.', '']:
                    current_equipment = equipment_number
                
                # Component found
                if current_equipment and component_name and component_name not in ['PARTS', '']:
                    if current_equipment in equipment_map:
                        # Find matching component in equipment
                        for comp_data in equipment_map[current_equipment].components:
                            if comp_data.component_name == component_name:
                                comp_data.row_index = current_row
                                logger.debug(f"Mapped {current_equipment}/{component_name} to row {current_row}")
                                break
                
                current_row += 1
            
            logger.info("✅ Component rows mapped")
        
        except Exception as e:
            logger.error(f"Error mapping rows: {str(e)}")
    
    def _fill_excel_data(self, ws, equipment_map: Dict[str, EquipmentData]):
        """Fill Excel data into template"""
        try:
            for equipment_data in equipment_map.values():
                for component_data in equipment_data.components:
                    if component_data.row_index:
                        row = component_data.row_index
                        
                        # Fill columns based on masterfile structure
                        ws[f'G{row}'] = component_data.data.get('fluid')
                        ws[f'H{row}'] = component_data.data.get('material_type')
                        ws[f'I{row}'] = component_data.data.get('spec')
                        ws[f'J{row}'] = component_data.data.get('grade')
                        ws[f'K{row}'] = component_data.data.get('insulation')
                        ws[f'L{row}'] = component_data.data.get('design_temp')
                        ws[f'M{row}'] = component_data.data.get('design_pressure')
                        ws[f'N{row}'] = component_data.data.get('operating_temp')
                        ws[f'O{row}'] = component_data.data.get('operating_pressure')
                        
                        logger.debug(f"Filled {equipment_data.equipment_number}/{component_data.component_name} at row {row}")
            
            logger.info("✅ Excel data filled")
        
        except Exception as e:
            logger.error(f"Error filling Excel: {str(e)}")
            raise
    
    def _get_cell_value(self, ws, cell_ref):
        """Safely get cell value"""
        try:
            return ws[cell_ref].value
        except:
            return None


# ============================================================================
# POWERPOINT GENERATION
# ============================================================================

class PowerPointReportGenerator:
    """Generate PowerPoint reports from extracted data"""
    
    # Rigid equipment sequence as per prototype requirements
    EQUIPMENT_SEQUENCE = [
        'V-001', 'V-002', 'V-003', 'V-004', 'V-005', 'V-006',
        'H-001', 'H-002', 'H-003', 'H-004'
    ]
    
    def __init__(self, template_path: str, log_callback=None):
        """Initialize with PowerPoint template path"""
        self.template_path = template_path
        if not os.path.exists(template_path):
            raise FileNotFoundError(f"PowerPoint template not found: {template_path}")
        
        self.log_callback = log_callback or (lambda msg: print(f"PPT: {msg}"))
        
        # Text box positions from Slide 0
        self.text_box_templates = {}
        
        # Font settings
        self.text_box_font = {
            'name': 'Arial',
            'size': Pt(10),
            'bold': False,
            'color': RGBColor(0, 0, 0)  # Black
        }
        
        self.table_font = {
            'name': 'Arial',
            'size': Pt(8),
            'bold': False,
            'color': RGBColor(0, 0, 0)  # Black
        }
    
    def log(self, message: str):
        """Log message using callback"""
        if self.log_callback:
            self.log_callback(message)
    
    def generate_from_equipment(self, equipment_list: List[Equipment]) -> bytes:
        """
        Generate PowerPoint report from equipment list.
        
        Args:
            equipment_list: List of Equipment objects from database
        
        Returns:
            PowerPoint file bytes
        """
        try:
            self.log(f"Generating PowerPoint report for {len(equipment_list)} equipment")
            
            # Load template
            prs = Presentation(self.template_path)
            
            # Extract text box positions from Slide 0
            self._extract_text_box_positions(prs.slides[0])
            
            # Build equipment data map
            equipment_map = {}
            for equipment in equipment_list:
                equip_data = EquipmentData(equipment)
                for component in equipment.components:
                    comp_data = ComponentData(component)
                    equip_data.components.append(comp_data)
                equipment_map[equipment.equipment_number] = equip_data
            
            # Fill each equipment's slide according to rigid sequence
            self._fill_slides_by_sequence(prs, equipment_map)
            
            # Save to bytes
            import io
            output = io.BytesIO()
            prs.save(output)
            output.seek(0)
            
            self.log("✅ PowerPoint report generated successfully")
            return output.getvalue()
        
        except Exception as e:
            self.log(f"❌ Error generating PowerPoint: {str(e)}")
            raise
    
    def _extract_text_box_positions(self, slide0):
        """Extract text box positions from Slide 0."""
        self.text_box_templates = {}
        
        for shape in slide0.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame.text:
                text = shape.text_frame.text.strip()
                if text in ["V-001", "Air Receiver", "MLK PMT 10101"]:
                    self.text_box_templates[text] = {
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'text': text
                    }
                    self.log(f"Found template text box: '{text}'")
    
    def _fill_slides_by_sequence(self, prs: Presentation, equipment_map: Dict[str, EquipmentData]):
        """Fill PowerPoint slides following the rigid equipment sequence"""
        try:
            total_slides = len(prs.slides)
            
            for slide_idx, expected_equip_no in enumerate(self.EQUIPMENT_SEQUENCE):
                # Check if we have this slide
                if slide_idx >= total_slides:
                    self.log(f"⚠️ No slide {slide_idx} for {expected_equip_no}")
                    break
                
                slide = prs.slides[slide_idx]
                
                # Check if we have this equipment
                if expected_equip_no in equipment_map:
                    equipment_data = equipment_map[expected_equip_no]
                    self.log(f"Filling slide {slide_idx} with {expected_equip_no}")
                    
                    # Add text boxes with Arial 10 font
                    self._add_text_boxes_to_slide(slide, equipment_data)
                    
                    # Fill equipment table with smart matching
                    self._fill_equipment_table(slide, equipment_data)
                    
                    self.log(f"✅ Filled slide {slide_idx} for {expected_equip_no}")
                else:
                    self.log(f"⚠️ Missing equipment {expected_equip_no} for slide {slide_idx}")
                    # Leave slide as-is (template with empty data)
        
        except Exception as e:
            self.log(f"Error filling slides: {str(e)}")
            raise
    
    def _add_text_boxes_to_slide(self, slide, equipment_data: EquipmentData):
        """Add text boxes to slide using template positions with Arial 10 font."""
        text_mapping = {
            "V-001": equipment_data.equipment_number,
            "Air Receiver": equipment_data.description or "",
            "MLK PMT 10101": equipment_data.pmt_number or f"PMT_{equipment_data.equipment_number}"
        }
        
        for template_text, equipment_value in text_mapping.items():
            if template_text in self.text_box_templates:
                self._create_text_box_with_font(slide, template_text, equipment_value)
    
    def _create_text_box_with_font(self, slide, template_text: str, equipment_value: str):
        """Create a text box on slide with Arial 10 font."""
        try:
            template = self.text_box_templates[template_text]
            
            textbox = slide.shapes.add_textbox(
                template['left'], 
                template['top'], 
                template['width'], 
                template['height']
            )
            
            text_frame = textbox.text_frame
            text_frame.clear()
            
            # Clear any default paragraphs
            while len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
                text_frame._element.remove(p._element)
            
            # Add new paragraph WITHOUT newline
            p = text_frame.add_paragraph()
            p.text = str(equipment_value)
            
            # Set paragraph alignment
            p.alignment = PP_ALIGN.LEFT
            
            # Set font to Arial 10
            run = p.runs[0]
            run.font.name = self.text_box_font['name']
            run.font.size = self.text_box_font['size']
            run.font.bold = self.text_box_font['bold']
            
            # Set font color
            run.font.color.rgb = self.text_box_font['color']
            
        except Exception as e:
            self.log(f"⚠️ Error creating text box '{template_text}': {e}")
    
    def _fill_equipment_table(self, slide, equipment_data: EquipmentData):
        """Fill equipment table with component data using smart matching"""
        try:
            # Find equipment table
            equipment_table = None
            for shape in slide.shapes:
                if hasattr(shape, 'table'):
                    table = shape.table
                    if len(table.columns) >= 8:
                        equipment_table = table
                        break
            
            if not equipment_table:
                self.log("No equipment table found on slide")
                return
            
            # Get expected component names from table (rows 2-4 in column 1)
            expected_components = []
            for row_idx in range(2, min(5, len(equipment_table.rows))):
                if row_idx < len(equipment_table.rows):
                    cell = equipment_table.cell(row_idx, 1)
                    if cell.text:
                        expected_name = cell.text.strip()
                        if expected_name and expected_name not in ['', 'COMPONENT', 'Part']:
                            expected_components.append((row_idx, expected_name))
            
            if not expected_components:
                self.log("No component names found in table")
                return
            
            # Fill each expected component using smart matching
            for row_idx, expected_name in expected_components:
                component_data = self._find_best_component_match(expected_name, equipment_data.components)
                if component_data:
                    self._fill_table_row(equipment_table, row_idx, component_data)
                    self.log(f"  ✓ Row {row_idx}: '{expected_name}' → '{component_data.component_name}'")
                else:
                    self.log(f"  ⚠️ No match for '{expected_name}'")
            
        except Exception as e:
            self.log(f"Warning: Error filling table: {str(e)}")
    
    def _find_best_component_match(self, expected_name: str, components: List[ComponentData]) -> Optional[ComponentData]:
        """Smart component matching logic"""
        if not components:
            return None
        
        expected_lower = expected_name.lower().strip()
        
        # 1. Try exact case-insensitive match
        for component in components:
            if component.component_name.lower().strip() == expected_lower:
                return component
        
        # 2. Try "contains" match (either direction)
        for component in components:
            comp_lower = component.component_name.lower().strip()
            if expected_lower in comp_lower or comp_lower in expected_lower:
                return component
        
        # 3. Try word overlap (split by non-alphanumeric)
        import re
        expected_words = set(re.findall(r'[a-z0-9]+', expected_lower))
        best_match = None
        best_score = 0
        
        for component in components:
            comp_words = set(re.findall(r'[a-z0-9]+', component.component_name.lower()))
            common_words = expected_words.intersection(comp_words)
            score = len(common_words)
            
            if score > best_score:
                best_score = score
                best_match = component
        
        if best_match and best_score > 0:
            return best_match
        
        # 4. Return first component if no better match found
        return components[0] if components else None
    
    def _fill_table_row(self, table, row_idx: int, component_data: ComponentData):
        """Fill a table row with component data using Arial 8 font WITHOUT newlines"""
        # Map component data to table columns
        column_mapping = {
            0: component_data.data.get('fluid', ''),      # Fluid
            2: component_data.data.get('spec', ''),       # Design Code/Spec
            3: component_data.data.get('material_type', ''),  # Material Type
            4: component_data.data.get('spec', ''),       # Spec
            5: component_data.data.get('grade', ''),      # Grade
            6: component_data.data.get('insulation', ''), # Insulation
            7: component_data.data.get('design_temp', ''),# Design Temp
            8: component_data.data.get('design_pressure', ''),  # Design Pressure
        }
        
        for col_idx, value in column_mapping.items():
            if col_idx < len(table.columns):
                # Remove any newlines from the value
                clean_value = str(value).replace('\n', '').replace('\r', '').strip()
                self._set_table_cell(table.cell(row_idx, col_idx), clean_value)
    
    def _set_table_cell(self, cell, value: str):
        """Set table cell value with Arial 8 WITHOUT newlines"""
        try:
            text_frame = cell.text_frame
            text_frame.clear()
            
            # Clear any existing paragraphs
            while len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
                text_frame._element.remove(p._element)
            
            # Add new paragraph WITHOUT newline
            p = text_frame.add_paragraph()
            p.text = value
            p.alignment = PP_ALIGN.CENTER
            
            if p.runs:
                run = p.runs[0]
                run.font.name = self.table_font['name']
                run.font.size = self.table_font['size']
                run.font.bold = self.table_font['bold']
                run.font.color.rgb = self.table_font['color']
        except Exception as e:
            self.log(f"Debug: Error setting cell: {str(e)}")


# ============================================================================
# MAIN REPORT GENERATION FUNCTIONS
# ============================================================================

async def generate_excel_report(
    db: Session,
    work_id: int,
    template_url: str
) -> str:
    """
    Generate Excel report from extracted equipment data.
    
    Args:
        db: Database session
        work_id: Work project ID
        template_url: Cloudinary URL of user's Excel template
    
    Returns:
        file_url (Cloudinary URL of generated report)
    """
    try:
        logger.info(f"Generating Excel report for work {work_id}")
        
        # ✓ FIXED: Download template with error handling
        logger.info(f"Downloading template from: {template_url}")
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(template_url)
            response.raise_for_status()  # ✓ Raise on 4xx/5xx
            
            template_bytes = response.content
        
        # ✓ FIXED: Validate template file
        if len(template_bytes) == 0:
            raise ValueError("Template file is empty - Cloudinary returned empty content")
        
        logger.info(f"Downloaded template: {len(template_bytes)} bytes")
        
        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.xlsx', delete=False) as tmp:
            tmp.write(template_bytes)
            template_path = tmp.name
        
        logger.debug(f"Template saved to: {template_path}")
        
        # Get equipment for this work
        equipment_list = db.query(Equipment).filter(
            Equipment.work_id == work_id
        ).all()
        
        if not equipment_list:
            raise ValueError("No equipment found for this work - cannot generate report")
        
        logger.info(f"Found {len(equipment_list)} equipment items")
        
        # Generate Excel
        generator = ExcelReportGenerator(template_path)
        excel_bytes = generator.generate_from_equipment(equipment_list)
        
        if len(excel_bytes) == 0:
            raise ValueError("Excel generation failed - no bytes produced")
        
        # Upload to Cloudinary
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        filename = f"work_{work_id}_excel_{timestamp}.xlsx"
        
        file_url = await upload_excel_to_cloudinary(
            file_bytes=excel_bytes,
            filename=filename
        )
        
        logger.info(f"[OK] Excel report uploaded: {file_url}")
        
        # Cleanup temp file
        os.unlink(template_path)
        
        return file_url
    
    except httpx.HTTPError as e:
        logger.error(f"[ERROR] Failed to download template from Cloudinary: {str(e)}")
        raise ValueError(f"Template download failed: {str(e)}")
    
    except ValueError as e:
        logger.error(f"[ERROR] Excel generation failed: {str(e)}")
        raise
    
    except Exception as e:
        logger.error(f"[ERROR] Unexpected error generating Excel: {str(e)}", exc_info=True)
        raise
    
    except Exception as e:
        logger.error(f"❌ Error generating Excel: {str(e)}")
        raise


async def generate_powerpoint_report(
    db: Session,
    work_id: int,
    template_url: str
) -> str:
    """
    Generate PowerPoint report from extracted equipment data.
    
    Args:
        db: Database session
        work_id: Work project ID
        template_url: Cloudinary URL of user's PowerPoint template
    
    Returns:
        file_url (Cloudinary URL of generated report)
    """
    try:
        logger.info(f"Generating PowerPoint report for work {work_id}")
        
        # ✓ FIXED: Download template with error handling
        logger.info(f"Downloading template from: {template_url}")
        async with httpx.AsyncClient(timeout=30.0) as client:
            response = await client.get(template_url)
            response.raise_for_status()  # ✓ Raise on 4xx/5xx
            
            template_bytes = response.content
        
        # ✓ FIXED: Validate template file
        if len(template_bytes) == 0:
            raise ValueError("Template file is empty - Cloudinary returned empty content")
        
        logger.info(f"Downloaded template: {len(template_bytes)} bytes")
        
        # Save to temp file
        with tempfile.NamedTemporaryFile(suffix='.pptx', delete=False) as tmp:
            tmp.write(template_bytes)
            template_path = tmp.name
        
        logger.debug(f"Template saved to: {template_path}")
        
        # Get equipment for this work
        equipment_list = db.query(Equipment).filter(
            Equipment.work_id == work_id
        ).all()
        
        if not equipment_list:
            raise ValueError("No equipment found for this work - cannot generate report")
        
        logger.info(f"Found {len(equipment_list)} equipment items")
        
        # Generate PowerPoint
        generator = PowerPointReportGenerator(template_path)
        ppt_bytes = generator.generate_from_equipment(equipment_list)
        
        if len(ppt_bytes) == 0:
            raise ValueError("PowerPoint generation failed - no bytes produced")
        
        # Upload to Cloudinary
        timestamp = datetime.utcnow().strftime("%Y%m%d_%H%M%S")
        filename = f"work_{work_id}_powerpoint_{timestamp}.pptx"
        
        file_url = await upload_ppt_to_cloudinary(
            file_bytes=ppt_bytes,
            filename=filename
        )
        
        logger.info(f"[OK] PowerPoint report uploaded: {file_url}")
        
        # Cleanup temp file
        os.unlink(template_path)
        
        return file_url
    
    except httpx.HTTPError as e:
        logger.error(f"[ERROR] Failed to download template from Cloudinary: {str(e)}")
        raise ValueError(f"Template download failed: {str(e)}")
    
    except ValueError as e:
        logger.error(f"[ERROR] PowerPoint generation failed: {str(e)}")
        raise
    
    except Exception as e:
        logger.error(f"[ERROR] Unexpected error generating PowerPoint: {str(e)}", exc_info=True)
        raise
