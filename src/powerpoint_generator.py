# powerpoint_generator.py
"""
PowerPoint Generator for AutoRBI - Custom Font Settings
"""
import os
import re
from typing import Dict, List, Optional
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from models import Equipment, Component


class PowerPointGenerator:
    def __init__(self, template_path: str, log_callback=None):
        if not template_path or not os.path.exists(template_path):
            raise FileNotFoundError(f"Template not found: {template_path}")
        
        self.template_path = template_path
        self.log_callback = log_callback or (lambda msg: print(f"PPT: {msg}"))
        
        # Text box positions from Slide 0
        self.text_box_templates = {}
        
        # Font settings
        self.text_box_font = {
            'name': 'Arial',
            'size': Pt(10),
            'bold': False,
            'color': RGBColor(0, 0, 0)  # Black
        }
        
        self.table_font = {
            'name': 'Arial',
            'size': Pt(8),
            'bold': False,
            'color': RGBColor(0, 0, 0)  # Black
        }
    
    def log(self, message: str) -> None:
        """Log a message using the callback."""
        if self.log_callback:
            self.log_callback(message)
    
    def generate_from_equipment_map(self, 
                                   equipment_map: Dict[str, Equipment], 
                                   output_path: str) -> bool:
        """Fill PowerPoint template with equipment data."""
        try:
            self.log(f"Loading template: {self.template_path}")
            
            # Load the template
            prs = Presentation(self.template_path)
            
            # Extract text box positions from Slide 0
            self._extract_text_box_positions(prs.slides[0])
            
            # Get equipment in sorted order, EXCLUDING V-001 since it's in template
            equipment_list = self._get_equipment_excluding_template(equipment_map)
            
            # Fill slides 1-9 with equipment data (skip slide 0 - it's the template)
            success = self._fill_slides_excluding_template(prs, equipment_list)
            
            if success:
                prs.save(output_path)
                self.log(f"✅ PowerPoint saved: {output_path}")
                return True
            else:
                self.log("⚠️ PowerPoint saved but some data may be missing")
                prs.save(output_path)
                return True
            
        except Exception as e:
            self.log(f"❌ Error generating PowerPoint: {str(e)}")
            import traceback
            traceback.print_exc()
            return False
    
    def _extract_text_box_positions(self, slide0):
        """Extract text box positions from Slide 0."""
        self.text_box_templates = {}
        
        for shape in slide0.shapes:
            if hasattr(shape, "text_frame") and shape.text_frame.text:
                text = shape.text_frame.text.strip()
                if text in ["V-001", "Air Receiver", "MLK PMT 10101"]:
                    self.text_box_templates[text] = {
                        'left': shape.left,
                        'top': shape.top,
                        'width': shape.width,
                        'height': shape.height,
                        'text': text
                    }
                    self.log(f"Found template text box: '{text}'")
    
    def _get_equipment_excluding_template(self, equipment_map: Dict[str, Equipment]) -> List[Equipment]:
        """Get equipment sorted by number, EXCLUDING V-001 (already in template)."""
        equipment_list = []
        
        for equip_no, equipment in equipment_map.items():
            # Skip V-001 since it's already in the template (Slide 0)
            if equip_no.upper() != "V-001":
                equipment_list.append(equipment)
                
        self.log(f"Selected {len(equipment_list)} equipment (excluding V-001 from template)")
        return equipment_list
    
    def _fill_slides_excluding_template(self, presentation, equipment_list: List[Equipment]) -> bool:
        """Fill slides 1-9 with equipment data (skip slide 0)."""
        total_slides = len(presentation.slides)
        
        # We have slides 1-9 to fill (10 slides total, 0-9)
        slides_to_fill = min(9, total_slides - 1)  # Skip slide 0
        
        self.log(f"Template has {total_slides} slides")
        self.log(f"Filling slides 1-{slides_to_fill} with {len(equipment_list)} equipment")
        
        for i in range(slides_to_fill):
            slide_idx = i + 1  # Start from slide 1 (skip slide 0)
            
            if i < len(equipment_list):
                equipment = equipment_list[i]
                slide = presentation.slides[slide_idx]
                
                self.log(f"  Slide {slide_idx} ({i+1}/{slides_to_fill}): {equipment.equipment_number}")
                
                # Add text boxes with Arial 10
                self._add_text_boxes_to_slide(slide, equipment)
                
                # Fill equipment table with Arial 8
                self._fill_equipment_table(slide, equipment)
            else:
                self.log(f"  Slide {slide_idx}: No more equipment available")
                # Leave slide as-is (template with empty data)
        
        return True
    
    def _add_text_boxes_to_slide(self, slide, equipment: Equipment):
        """Add text boxes to slide using template positions with Arial 10 font."""
        text_mapping = {
            "V-001": equipment.equipment_number,
            "Air Receiver": equipment.equipment_description,
            "MLK PMT 10101": equipment.pmt_number or f"PMT_{equipment.equipment_number}"
        }
        
        for template_text, equipment_value in text_mapping.items():
            if template_text in self.text_box_templates:
                self._create_text_box_with_font(slide, template_text, equipment_value)
    
    def _create_text_box_with_font(self, slide, template_text: str, equipment_value: str):
        """Create a text box on slide with Arial 10 font."""
        try:
            template = self.text_box_templates[template_text]
            
            textbox = slide.shapes.add_textbox(
                template['left'], 
                template['top'], 
                template['width'], 
                template['height']
            )
            
            text_frame = textbox.text_frame
            text_frame.clear()
            
            # Clear any default paragraphs
            while len(text_frame.paragraphs) > 0:
                p = text_frame.paragraphs[0]
                text_frame._element.remove(p._element)
            
            # Add new paragraph
            p = text_frame.add_paragraph()
            p.text = str(equipment_value)
            
            # Set paragraph alignment
            p.alignment = PP_ALIGN.LEFT
            
            # Set font to Arial 10
            run = p.runs[0]
            run.font.name = self.text_box_font['name']
            run.font.size = self.text_box_font['size']
            run.font.bold = self.text_box_font['bold']
            
            # Set font color
            run.font.color.rgb = self.text_box_font['color']
            
        except Exception as e:
            self.log(f"⚠️ Error creating text box '{template_text}': {e}")
    
    def _fill_equipment_table(self, slide, equipment: Equipment):
        """Fill the equipment table on slide with Arial 8 font."""
        try:
            # Find equipment table (first table with enough columns)
            equipment_table = None
            for shape in slide.shapes:
                if hasattr(shape, "table"):
                    table = shape.table
                    if len(table.columns) >= 8:  # Equipment table has many columns
                        equipment_table = table
                        break
            
            if not equipment_table:
                self.log(f"    No equipment table found")
                return
            
            # Get component names from table (rows 2-4 in column 1)
            expected_components = []
            for row_idx in range(2, min(5, len(equipment_table.rows))):
                if row_idx < len(equipment_table.rows) and 1 < len(equipment_table.columns):
                    cell_text = equipment_table.cell(row_idx, 1).text.strip()
                    if cell_text:
                        expected_components.append((row_idx, cell_text))
            
            if not expected_components:
                self.log(f"    No component names found in table")
                return
            
            # Fill each component row
            for row_idx, expected_name in expected_components:
                component = self._find_best_component_match(expected_name, equipment.components)
                if component:
                    self._fill_table_row_with_font(equipment_table, row_idx, component.existing_data)
                    self.log(f"    ✓ Row {row_idx} ('{expected_name}'): {component.component_name}")
                else:
                    self.log(f"    ⚠️ No component match for '{expected_name}'")
            
        except Exception as e:
            self.log(f"    Error filling table: {e}")
    
    def _find_best_component_match(self, expected_name: str, components: List[Component]) -> Optional[Component]:
        """Find the best matching component for the expected name."""
        if not components:
            return None
        
        expected_lower = expected_name.lower()
        
        # Try exact match first
        for component in components:
            if component.component_name.lower() == expected_lower:
                return component
        
        # Try contains match
        for component in components:
            comp_lower = component.component_name.lower()
            if expected_lower in comp_lower or comp_lower in expected_lower:
                return component
        
        # Try word overlap
        expected_words = set(re.findall(r'\w+', expected_lower))
        for component in components:
            comp_words = set(re.findall(r'\w+', component.component_name.lower()))
            if expected_words.intersection(comp_words):
                return component
        
        # Return first component if no match
        return components[0]
    
    def _fill_table_row_with_font(self, table, row_idx: int, comp_data: dict):
        """Fill a table row with component data using Arial 8 font."""
        # Column mapping
        column_mapping = {
            'fluid': 0,
            'design_code': 2,
            'material_type': 3,
            'spec': 4,
            'grade': 5,
            'insulation': 6,
            'design_temp': 7,
            'design_pressure': 8,
            'risk_rating': 9,
        }
        
        for data_key, col_idx in column_mapping.items():
            if col_idx < len(table.columns):
                value = comp_data.get(data_key, "")
                
                if value:
                    cell = table.cell(row_idx, col_idx)
                    
                    # Only fill if cell is empty or we have important data
                    current_text = cell.text.strip()
                    should_fill = (not current_text or 
                                  data_key in ['fluid', 'design_code', 'material_type'])
                    
                    if should_fill:
                        # Clear existing text
                        cell.text = ""
                        
                        # Clear existing paragraphs
                        text_frame = cell.text_frame
                        while len(text_frame.paragraphs) > 0:
                            p = text_frame.paragraphs[0]
                            text_frame._element.remove(p._element)
                        
                        # Add new paragraph with Arial 8
                        p = text_frame.add_paragraph()
                        p.text = str(value)
                        p.alignment = PP_ALIGN.CENTER
                        
                        # Set font to Arial 8
                        if p.runs:
                            run = p.runs[0]
                            run.font.name = self.table_font['name']
                            run.font.size = self.table_font['size']
                            run.font.bold = self.table_font['bold']
                            run.font.color.rgb = self.table_font['color']